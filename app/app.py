from flask import Flask, render_template, request
import mysql.connector, random
import openai
import config as c
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from pptx import Presentation

openai.organization = c.openai_org_id
openai.api_key = c.openai_secret_key

app = Flask(__name__)
app.secret_key = c.flask_secret_key

DATABASE_TABLE = 'question'
db = mysql.connector.connect(
    host = c.mysql_hostname,
    user = c.mysql_username,
    password = c.mysql_password,
    database = c.mysql_database
)
cursor = db.cursor()

TYPES = {
    'Movies' : 1,
    'TV' : 2,
    'Books' : 3, 
    'Games' : 4,
}
GENRES = {
    'Action' : 1,
    'Comedy' : 2,
    'Horror' : 3,
    'Romance' : 4,
    'Fantasy' : 5
}

DIFFICULTIES = {
    'Easy' : 1,
    'Medium' : 2,
    'Hard' : 3
}

PDF_STYLE = {
    'content': {
        'fontName': 'Helvetica-Bold',
        'fontSize': 12,
        'textColor': colors.blue,
        'alignment': 1,  # Center alignment
        'spaceAfter': 10,
    },
    'answer': {
        'fontName': 'Helvetica',
        'fontSize': 10,
        'textColor': colors.black,
        'alignment': 0,  # Left alignment
        'spaceAfter': 5,
    }
}

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/submit', methods=['POST'])
def submit():
    num_questions = int(request.form.get('number'))
    type = request.form.get('type')
    genre = request.form.get('genre')
    difficulty = request.form.get('difficulty')

    questions = get_questions(num_questions, TYPES[type], GENRES[genre], DIFFICULTIES[difficulty])

    export_pdf(questions)
    export_ppt(questions)

    return "Form submitted successfully!"

def get_questions(n, type, genre, difficulty):

    cursor.execute(
        f"SELECT * FROM {DATABASE_TABLE} WHERE Genre = %s AND Type = %s AND Difficulty = %s",
        (type, genre, difficulty)
    )

    return [{'content': row[4], 'answer': row[5]} for row in random.sample(cursor.fetchall(), n)]

def export_pdf(questions):
    doc = SimpleDocTemplate('TigerHacksTrivia.pdf', pagesize=letter)
    styles = getSampleStyleSheet()
    elements = []

    for question in questions:
        content = question['content']
        answer = question['answer']

        # Add question
        elements.append(Paragraph(f'<font name="{PDF_STYLE["content"]["fontName"]}" size="{PDF_STYLE["content"]["fontSize"]}" color="{PDF_STYLE["content"]["textColor"]}">{content}</font>', styles['Normal']))
        elements.append(Spacer(1, PDF_STYLE["content"]["spaceAfter"]))

        # Add answer
        elements.append(Paragraph(f'<font name="{PDF_STYLE["answer"]["fontName"]}" size="{PDF_STYLE["answer"]["fontSize"]}" color="{PDF_STYLE["answer"]["textColor"]}">{answer}</font>', styles['Normal']))
        elements.append(Spacer(1, PDF_STYLE["answer"]["spaceAfter"]))

    doc.build(elements)

def export_ppt(questions):
  prs = Presentation()
  
  for question in questions:
    content = question['content']
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"{content}"

  for question in questions:
    content = question['content']
    answer = question['answer']

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"{content}"
    slide.placeholders[1].text = f"{answer}"

  prs.save('TigerHacksTrivia.pptx')

def gpt_api():
    pass
    # prompt = "tell me a joke"

    # # Make the API request
    # response = openai.Completion.create(
    #     engine="davinci",
    #     prompt=prompt,
    #     max_tokens=50,  # Adjust max tokens as needed
    # )

    # print(response)
    # print()
    # joke = response.choices[0].text.strip()
    # print(joke)

if __name__ == '__main__':
    app.run(debug=True)






